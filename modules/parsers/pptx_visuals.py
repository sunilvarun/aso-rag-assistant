# modules/parsers/pptx_visuals.py
from __future__ import annotations
import re
from dataclasses import dataclass
from typing import List, Tuple, Dict, Optional
from datetime import date

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------
# Data classes
# ---------------------
@dataclass
class BBox:
    x: float; y: float; w: float; h: float
    @property
    def cx(self): return self.x + self.w / 2
    @property
    def cy(self): return self.y + self.h / 2

@dataclass
class TextBox:
    slide: int
    text: str
    bbox: BBox
    font_size: float
    bold: bool
    group_id: Optional[int] = None

@dataclass
class ShapeBox:
    slide: int
    kind: str               # "line" | "circle" | "rect" | "other"
    bbox: BBox
    group_id: Optional[int] = None

@dataclass
class Axis:
    slide: int
    y0: float; y1: float
    x0: float; x1: float

@dataclass
class DateTok:
    slide: int
    raw: str
    iso: Optional[str]
    bbox: BBox

@dataclass
class Milestone:
    slide: int
    title: str
    date_iso: Optional[str]
    date_raw: Optional[str]
    title_bbox: BBox
    date_bbox: Optional[BBox]
    confidence: float
    source: str

@dataclass
class SpanRow:
    slide: int
    title: str
    start_iso: Optional[str]
    end_iso: Optional[str]
    start_raw: Optional[str]
    end_raw: Optional[str]
    title_bbox: BBox
    confidence: float
    source: str

# ---------------------
# Regex & helpers
# ---------------------
MONTHS = {"jan":1,"feb":2,"mar":3,"apr":4,"may":5,"jun":6,"jul":7,"aug":8,"sep":9,"oct":10,"nov":11,"dec":12}

DATE_RE = re.compile(
    r"\b("
    r"Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|"
    r"Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:t)?(?:ember)?|"
    r"Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)"
    r"\.?\s+(\d{1,2})\b",
    re.I,
)
DATE_RE_DMY = re.compile(
    r"\b(\d{1,2})\s+("
    r"Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|"
    r"Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:t)?(?:ember)?|"
    r"Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)\b",
    re.I,
)
RANGE_RE = re.compile(rf"({DATE_RE.pattern})\s*(?:–|-|to)\s*({DATE_RE.pattern})", re.I)

def _ppt_len(x) -> float: return float(x)
def _looks_like_date(s: str) -> bool:
    """Return True if the string looks like a date (e.g., 'May 17', '17 May')."""
    return bool(DATE_RE.search(s) or DATE_RE_DMY.search(s))
def _shape_bbox(sh) -> BBox: return BBox(_ppt_len(sh.left), _ppt_len(sh.top), _ppt_len(sh.width), _ppt_len(sh.height))
def _percentile(values: List[float], p: float) -> float:
    if not values: return 0.0
    vs = sorted(values); idx = max(0, min(len(vs)-1, int(round(p*(len(vs)-1)))))
    return vs[idx]

# ---------------------
# Extract shapes & text
# ---------------------
def _text_boxes(slide_idx: int, slide) -> List[TextBox]:
    out: List[TextBox] = []
    for sh in slide.shapes:
        if not hasattr(sh, "has_text_frame") or not sh.has_text_frame:
            continue
        text = (sh.text or "").strip()
        text = text.replace("\xa0"," ").replace("\n"," ").replace("\r"," ")
        if not text:
            continue
        fs=None; bold=False
        try:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size: fs=float(r.font.size.pt)
                    if r.font.bold: bold=True
                    break
                if fs: break
        except Exception:
            pass
        fs = fs or 12.0
        out.append(TextBox(slide_idx, text, _shape_bbox(sh), fs, bold))
    return out

def _shape_boxes(slide_idx: int, slide) -> List[ShapeBox]:
    out: List[ShapeBox] = []
    def walk(container):
        for sh in container.shapes:
            st = sh.shape_type
            if st == MSO_SHAPE_TYPE.GROUP:
                walk(sh); continue
            kind = "other"
            w, h = float(sh.width), float(sh.height)
            if st == MSO_SHAPE_TYPE.LINE or (h > 0 and (w/h) > 10): kind = "line"
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE and h>0 and abs(w-h) <= max(12.0, 0.25*h): kind = "circle"
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE: kind = "rect"
            out.append(ShapeBox(slide_idx, kind, _shape_bbox(sh)))
    walk(slide)
    return out

# ---------------------
# Axis & alignment
# ---------------------
def _detect_axis(shapes: List[ShapeBox]) -> Optional[Axis]:
    lines = [s for s in shapes if s.kind == "line"]
    best  = max(lines, key=lambda s: s.bbox.w, default=None)
    if best:
        y = best.bbox.cy
        return Axis(best.slide, y-8, y+8, best.bbox.x, best.bbox.x+best.bbox.w)
    circles = [s for s in shapes if s.kind == "circle"]
    if len(circles) >= 3:
        ys = [c.bbox.cy for c in circles]
        y_min, y_max = _percentile(ys, 0.25), _percentile(ys, 0.75)
        row = [c for c in circles if (y_min-12) <= c.bbox.cy <= (y_max+12)]
        if row:
            xs = [c.bbox.cx for c in row]
            return Axis(row[0].slide, y_min-8, y_max+8, min(xs), max(xs))
    return None

def _axis_from_dates(toks: List[DateTok], slide_h: float) -> Optional[Axis]:
    if len(toks) < 2: return None
    ys = [d.bbox.cy for d in toks]; xs=[d.bbox.cx for d in toks]
    y_med = _percentile(ys, 0.5)
    band = 0.04 * slide_h  # 4% of slide height
    return Axis(toks[0].slide, y_med - band/2, y_med + band/2, min(xs), max(xs))

def _in_band(y: float, axis: Axis, pad: float = 0) -> bool:
    return (axis.y0 - pad) <= y <= (axis.y1 + pad)

# ---------------------
# Dates & titles
# ---------------------
def _extract_year_headers(texts: List[TextBox]) -> List[TextBox]:
    """Pick likely year headers: 4-digit numbers, prefer larger fonts."""
    cands = [t for t in texts if (len(t.text.strip()) == 4 and t.text.strip().isdigit())]
    cutoff = _percentile([t.font_size for t in cands], 0.7) if cands else 0
    return [t for t in cands if t.font_size >= cutoff]

def _parse_mmdd(raw: str, year: Optional[int]) -> Optional[str]:
    m = DATE_RE.search(raw) or DATE_RE_DMY.search(raw)
    if not m: return None
    if m.re is DATE_RE:
        mon = MONTHS[m.group(1).lower()[:3]]; day = int(m.group(2))
    else:
        day = int(m.group(1)); mon = MONTHS[m.group(2).lower()[:3]]
    y = year or date.today().year
    return f"{y:04d}-{mon:02d}-{day:02d}"

def _assign_years(tokens: List[DateTok], year_headers: List[TextBox]) -> None:
    headers_sorted = sorted(year_headers, key=lambda t: t.bbox.cx)
    for tok in tokens:
        year=None; best_dx=1e12
        for h in headers_sorted:
            dx = abs(tok.bbox.cx - h.bbox.cx)
            if dx < best_dx:
                year=int(h.text.strip()); best_dx=dx
        tok.iso = _parse_mmdd(tok.raw, year)

def _detect_dates_inline(texts: List[TextBox]) -> List[DateTok]:
    toks: List[DateTok] = []
    for t in texts:
        for m in DATE_RE.finditer(t.text):
            toks.append(DateTok(t.slide, m.group(0), None, t.bbox))
        for m in DATE_RE_DMY.finditer(t.text):
            toks.append(DateTok(t.slide, f"{m.group(2)} {m.group(1)}", None, t.bbox))
    return toks

def _detect_dates_split_boxes(texts: List[TextBox], axis: Axis, year_headers: List[TextBox], slide_h: float) -> List[DateTok]:
    mon_only = re.compile(r"^(Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:t)?(?:ember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)\.?$", re.I)
    day_only = re.compile(r"^\d{1,2}$")
    months, days = [], []
    for t in texts:
        s = t.text.strip()
        if mon_only.fullmatch(s): months.append(t)
        elif day_only.fullmatch(s) and 1 <= int(s) <= 31: days.append(t)

    if not months or not days: return []

    toks: List[DateTok] = []
    used = set()
    for mb in months:
        best=None; best_sc=1e18
        for db in days:
            if db in used: continue
            dx = abs(mb.bbox.cx - db.bbox.cx)
            dy = abs(mb.bbox.cy - db.bbox.cy)
            sc = dx + dy
            if sc < best_sc:
                best_sc = sc; best = db
        if not best: continue
        used.add(best)
        raw = f"{mb.text.strip()} {best.text.strip()}"
        cx = (mb.bbox.cx + best.bbox.cx)/2; cy = (mb.bbox.cy + best.bbox.cy)/2
        toks.append(DateTok(mb.slide, raw, None, BBox(cx-1, cy-1, 2, 2)))

    _assign_years(toks, year_headers)
    pad = 0.30 * slide_h
    return [d for d in toks if _in_band(d.bbox.cy, axis, pad=pad)]

def _detect_dates(texts: List[TextBox], axis: Axis, year_headers: List[TextBox], slide_h: float) -> List[DateTok]:
    inline = _detect_dates_inline(texts)
    _assign_years(inline, year_headers)
    pad = 0.30 * slide_h   # generous vertical band around axis
    near = [d for d in inline if _in_band(d.bbox.cy, axis, pad=pad)]
    if len(near) < 2:
        near += _detect_dates_split_boxes(texts, axis, year_headers, slide_h)
        near.sort(key=lambda d: d.bbox.cx)
        # dedupe by ~1% of axis width (or slide width fallback)
        dedup: List[DateTok] = []
        x_eps = 0.01 * (axis.x1 - axis.x0 if axis.x1 > axis.x0 else 1.0)
        for d in near:
            if not dedup or abs(d.bbox.cx - dedup[-1].bbox.cx) > x_eps:
                dedup.append(d)
        return dedup
    return near

def _detect_titles(texts: List[TextBox], axis: Axis, slide_h: float) -> List[TextBox]:
    mid_y = (axis.y0 + axis.y1) / 2
    near = [t for t in texts if abs(t.bbox.cy - mid_y) <= 0.50 * slide_h]  # wide window
    month_only = re.compile(r"^(Today|Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:t)?(?:ember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)$", re.I)
    out: List[TextBox] = []
    for t in near:
        s = t.text.strip()
        if month_only.fullmatch(s): continue
        if not re.search(r"[A-Za-z]", s): continue
        out.append(t)
    return out

# ---------------------
# Pairing
# ---------------------

def _pair_axisless(
    dates: List[DateTok],
    titles: List[TextBox],
    axis: Axis,
    slide_w: float,
    slide_h: float,
) -> List[Milestone]:
    """
    Pair each date token with the best nearby NON-DATE title box using x & y proximity.
    This fixes cases where titles are above/below the month-day label.
    """
    out: List[Milestone] = []
    if not dates or not titles:
        return out

    # keep only titles that are not themselves dates (e.g., "May 17")
    cand_titles = [t for t in titles if not _looks_like_date(t.text)]

    if not cand_titles:
        return out

    mid_y = (axis.y0 + axis.y1) / 2

    # tolerances relative to slide size (EMU-safe)
    x_tol_primary   = 0.22 * slide_w
    x_tol_secondary = 0.30 * slide_w
    y_tol_primary   = 0.18 * slide_h
    y_tol_secondary = 0.30 * slide_h

    def best_title_for_date(d: DateTok) -> Optional[TextBox]:
        def pick(pool, x_tol, y_tol):
            cands: List[Tuple[float, TextBox]] = []
            for t in pool:
                dx = abs(t.bbox.cx - d.bbox.cx)
                dy = abs(t.bbox.cy - d.bbox.cy)
                if dx <= x_tol and dy <= y_tol:
                    # Score favors x alignment, then closeness to axis (smaller |cy - mid_y|), plus mild font boost
                    axis_prox = 1 - min(abs(t.bbox.cy - mid_y) / (0.5 * slide_h), 1.0)
                    score = (1 - dx / x_tol) * 0.6 + (1 - dy / y_tol) * 0.3 + (t.font_size / 24.0) * 0.1
                    cands.append((score, t))
            if not cands:
                return None
            return max(cands, key=lambda x: x[0])[1]

        # Try tighter window first, then expand
        t = pick(cand_titles, x_tol_primary, y_tol_primary)
        if t is None:
            t = pick(cand_titles, x_tol_secondary, y_tol_secondary)
        return t

    for d in dates:
        t = best_title_for_date(d)
        if t:
            # confidence primarily from x alignment (stable heuristic)
            conf_x = 1 - min(abs(t.bbox.cx - d.bbox.cx) / x_tol_secondary, 1.0)
            conf_y = 1 - min(abs(t.bbox.cy - d.bbox.cy) / y_tol_secondary, 1.0)
            conf   = 0.7 * conf_x + 0.3 * conf_y
            out.append(Milestone(
                slide=d.slide,
                title=t.text.strip(),
                date_iso=d.iso,
                date_raw=d.raw,
                title_bbox=t.bbox,
                date_bbox=d.bbox,
                confidence=conf,
                source=""
            ))
    return out

# ---------------------
# Spans
# ---------------------
def _detect_spans_from_titles(titles: List[TextBox], year_headers: List[TextBox], source: str, slide_idx: int) -> List[SpanRow]:
    rows: List[SpanRow] = []
    for t in titles:
        m = RANGE_RE.search(t.text)
        if not m: continue
        left_raw, right_raw = m.group(1), m.group(4)
        # naive year attach: use the first year header (good enough for hackathon decks)
        y = int(year_headers[0].text.strip()) if year_headers else None
        start_iso = _parse_mmdd(left_raw, y); end_iso = _parse_mmdd(right_raw, y)
        rows.append(SpanRow(slide_idx, t.text.strip(), start_iso, end_iso, left_raw, right_raw, t.bbox, 0.9, source))
    return rows

# ---------------------
# Public entry
# ---------------------
def parse_pptx_visuals(pptx_path: str, status_legend: Dict[str, str] | None = None):
    prs = Presentation(pptx_path)
    slide_w_total = float(prs.slide_width)
    slide_h_total = float(prs.slide_height)

    captions: List[str] = []
    milestones: List[Milestone] = []
    spans: List[SpanRow] = []

    for i, slide in enumerate(prs.slides, start=1):
        shapes = _shape_boxes(i, slide)
        texts  = _text_boxes(i, slide)

        # Dates (inline) just to help axis fallback
        inline_dates = _detect_dates_inline(texts)

        # Axis: prefer shape-based, fallback to dates
        axis = _detect_axis(shapes)
        if not axis and inline_dates:
            axis = _axis_from_dates(inline_dates, slide_h_total)

        if not axis:
            raw = " ".join([t.text for t in texts])
            captions.append(f"Slide {i}: {raw[:1200]}")
            continue

        # Unified detection & pairing
        year_headers = _extract_year_headers(texts)
        date_toks    = _detect_dates(texts, axis, year_headers, slide_h_total)
        title_blks   = _detect_titles(texts, axis, slide_h_total)

        ms = _pair_axisless(date_toks, title_blks, axis, slide_w_total, slide_h_total)
        for m in ms: m.source = pptx_path
        milestones.extend(ms)

        spans.extend(_detect_spans_from_titles(title_blks, year_headers, pptx_path, i))

        # Caption for RAG
        vis_bits = [f"{m.title} ({m.date_raw})" for m in ms if m.title and m.date_raw]
        if vis_bits:
            captions.append(f"Slide {i}: " + "; ".join(vis_bits))
        else:
            raw = " ".join([t.text for t in texts])
            captions.append(f"Slide {i}: {raw[:1200]}")

    structured = {
        "milestones": [
            {"slide": m.slide, "title": m.title, "date": m.date_iso, "raw_date": m.date_raw,
             "confidence": round(m.confidence, 3), "source": m.source}
            for m in milestones if m.title
        ],
        "spans": [
            {"slide": s.slide, "title": s.title, "start_date": s.start_iso, "end_date": s.end_iso,
             "raw_left": s.start_raw, "raw_right": s.end_raw, "confidence": round(s.confidence, 3),
             "source": s.source}
            for s in spans if s.title
        ],
    }
    return captions, structured
