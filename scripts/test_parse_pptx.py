# scripts/test_parse_pptx.py
from modules.parsers.pptx_visuals import parse_pptx_visuals

# point this to one of your timeline decks
pptx_path = "/Users/sunilvarun/Documents/Python/my_corpus/Timeline - Multiple Complex.pptx"
#pptx_path = "/Users/sunilvarun/Documents/Python/my_corpus/Timeline - Dates Above Milestones below.pptx"
#pptx_path = "/Users/sunilvarun/Documents/Python/my_corpus/Timeline - Simple but overlapping text.pptx"
# --- Sanity check: what does the regex see? ---
from pptx import Presentation
from modules.parsers.pptx_visuals import DATE_RE, DATE_RE_DMY

prs = Presentation(pptx_path)
slide = prs.slides[0]
texts = []
for sh in slide.shapes:
    if hasattr(sh, "has_text_frame") and sh.has_text_frame and sh.has_text_frame:
        t = (sh.text or "")
        # mirror the parser's normalization
        t = t.replace("\xa0", " ").replace("\n", " ").replace("\r", " ")
        texts.append(t)

hits = []
for t in texts:
    if DATE_RE.search(t) or DATE_RE_DMY.search(t):
        hits.append(t)

print(f"[SANITY] Text boxes: {len(texts)} | date hits: {len(hits)}")
for h in hits[:8]:
    print("[HIT]", h)
print("-" * 60)


caps, stru = parse_pptx_visuals(pptx_path)
print(f"\nFound {len(stru['milestones'])} milestones, {len(stru['spans'])} spans")

print("---- Captions (for RAG) ----")
for c in caps:
    print(c)

print("\n---- Milestones ----")
for m in stru["milestones"]:
    print(m)

print("\n---- Spans ----")
for s in stru["spans"]:
    print(s)
    
print(f"\nFound {len(stru['milestones'])} milestones, {len(stru['spans'])} spans")
